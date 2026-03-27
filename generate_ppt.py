from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0A, 0x0B, 0x10)   # near-black
CARD    = RGBColor(0x12, 0x14, 0x20)   # dark card
ACCENT  = RGBColor(0x00, 0xF2, 0xFF)   # cyan
ACCENT2 = RGBColor(0xFF, 0x33, 0x66)   # red-pink
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0xAA, 0xAA, 0xAA)
GREEN   = RGBColor(0x00, 0xFF, 0xCC)
YELLOW  = RGBColor(0xFF, 0xD7, 0x00)

SLIDE_LAYOUT = prs.slide_layouts[6]   # blank

# ── Helpers ──────────────────────────────────────────────────────────────────
def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, color=CARD, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def accent_bar(slide, t=0.0, h=0.09, color=ACCENT):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(t), prs.slide_width, Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def badge(slide, text, l, t, w=1.2, h=0.32, bg_color=ACCENT, fg_color=BG):
    b = box(slide, l, t, w, h, bg_color)
    b.line.fill.background()
    tf = b.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = fg_color

def section_header(slide, label, color=ACCENT):
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(0.18), Inches(1.6), Inches(0.38))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = BG

def bullet_card(slide, l, t, w, h, title, points, title_color=ACCENT, pt_size=13):
    box(slide, l, t, w, h, CARD)
    txt(slide, title, l+0.18, t+0.12, w-0.26, 0.35,
        size=14, bold=True, color=title_color)
    yy = t + 0.48
    for pt in points:
        txt(slide, f"▸  {pt}", l+0.18, yy, w-0.28, 0.28,
            size=pt_size, color=WHITE)
        yy += 0.3

def divider(slide, t, color=ACCENT):
    s = slide.shapes.add_shape(1, Inches(0.5), Inches(t),
                                Inches(12.33), Inches(0.03))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(SLIDE_LAYOUT)
bg(sl)

# Top accent
accent_bar(sl, 0.0, 0.09, ACCENT)
accent_bar(sl, 7.41, 0.09, ACCENT2)

# Central glow box
box(sl, 1.5, 1.7, 10.33, 4.1, RGBColor(0x10, 0x12, 0x22))

txt(sl, "ShadowNet Nexus", 2.0, 1.9, 9.5, 1.2,
    size=54, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(sl, "v5.0", 2.0, 2.95, 9.5, 0.8,
    size=40, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
txt(sl, "Unified Live Forensic Intelligence System",
    2.0, 3.75, 9.5, 0.55,
    size=20, bold=False, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

divider(sl, 4.45)

txt(sl, "AI-Powered  ·  Real-Time  ·  Kernel-Level  ·  Self-Healing",
    2.0, 4.62, 9.5, 0.45,
    size=15, color=GREY, align=PP_ALIGN.CENTER)

txt(sl, "Built with Google Gemini 2.5 Flash + Raw NTFS Parsing + Behavioral GMM Analysis",
    1.8, 5.25, 9.8, 0.45,
    size=13, color=GREY, align=PP_ALIGN.CENTER, italic=True)

# Corner badges
badge(sl, "V5.0", 11.3, 6.9, 1.5, 0.35, ACCENT2, WHITE)
badge(sl, "Python", 0.35, 6.9, 1.2, 0.35, CARD, ACCENT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(SLIDE_LAYOUT)
bg(sl)
accent_bar(sl, 0.0, 0.07, ACCENT2)
section_header(sl, "THE PROBLEM", ACCENT2)

txt(sl, "Why Traditional Tools Fail Against Modern Attackers",
    0.5, 0.62, 12.5, 0.55, size=26, bold=True, color=WHITE)
divider(sl, 1.28, ACCENT2)

problems = [
    ("Living off the Land",
     "Attackers use built-in Windows tools (PowerShell, wevtutil, certutil).\nNo malicious file is ever written — standard AV has nothing to scan."),
    ("Timestomping",
     "Malware backdates its own file creation date to appear as a\nlegitimate old system file, defeating timeline analysis."),
    ("Volume Shadow Deletion",
     "Ransomware silently runs 'vssadmin delete shadows' to remove\nall Windows backups before encrypting your data."),
    ("High-Entropy Wiping",
     "Data wipers overwrite files with pure random bytes or all-zeros.\nThe original data is gone before anyone notices."),
    ("Automated Bot Input",
     "Scripts steal passwords with keystroke speeds physically\nimpossible for humans — standard tools cannot tell the difference."),
    ("Log Destruction",
     "Skilled attackers clear Windows Event Logs and NTFS change journals,\nerasing all records of what happened on the system."),
]

cols = [(0.4, 1.45), (4.55, 1.45), (8.7, 1.45),
        (0.4, 3.85),  (4.55, 3.85), (8.7, 3.85)]

for i, ((l, t), (title, desc)) in enumerate(zip(cols, problems)):
    box(sl, l, t, 3.9, 2.1, CARD)
    # Red dot
    badge(sl, "✕", l+0.1, t+0.1, 0.38, 0.32, ACCENT2, WHITE)
    txt(sl, title, l+0.58, t+0.12, 3.2, 0.36, size=14, bold=True, color=ACCENT2)
    txt(sl, desc, l+0.15, t+0.55, 3.6, 1.4, size=11.5, color=GREY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What ShadowNet Does
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(SLIDE_LAYOUT)
bg(sl)
accent_bar(sl, 0.0, 0.07, ACCENT)
section_header(sl, "THE SOLUTION", ACCENT)

txt(sl, "ShadowNet Nexus v5.0 — How It Thinks Differently",
    0.5, 0.62, 12.5, 0.55, size=26, bold=True, color=WHITE)
divider(sl, 1.28, ACCENT)

txt(sl,
    "ShadowNet does NOT scan files for known virus signatures.\n"
    "It observes the BEHAVIOR of the operating system itself — at the hardware level.",
    0.5, 1.4, 12.3, 0.65, size=15, color=GREY, italic=True)

solutions = [
    (ACCENT,  "Reads Raw Disk Metal",
     "Opens the hard drive the same way the OS kernel does — bypassing all software layers.\n"
     "Rootkits and malware literally cannot hide from this."),
    (GREEN,   "Measures Data Randomness",
     "Computes the mathematical entropy of every cluster on the disk.\n"
     "Ransomware and wipers leave a unique statistical fingerprint that cannot be faked."),
    (YELLOW,  "Analyzes Human Heartbeat",
     "Measures the exact nanosecond timing between every keystroke.\n"
     "Automated bots type at speeds no human can physically match."),
    (ACCENT2, "Understands Command Intent",
     "Sends every detected command to Google Gemini AI.\n"
     "The AI reads the intent — not just the syntax — and produces a full threat verdict."),
]

xx = 0.4
for color, title, desc in solutions:
    box(sl, xx, 2.2, 2.95, 4.5, CARD)
    s = sl.shapes.add_shape(1, Inches(xx), Inches(2.2), Inches(2.95), Inches(0.18))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    txt(sl, title, xx+0.15, 2.45, 2.65, 0.42, size=14, bold=True, color=color)
    txt(sl, desc,  xx+0.15, 2.95, 2.65, 3.5,  size=12, color=WHITE)
    xx += 3.15

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — System Architecture
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(SLIDE_LAYOUT)
bg(sl)
accent_bar(sl, 0.0, 0.07, ACCENT)
section_header(sl, "ARCHITECTURE", ACCENT)

txt(sl, "7-Layer System Architecture",
    0.5, 0.62, 12.5, 0.55, size=26, bold=True, color=WHITE)
divider(sl, 1.28, ACCENT)

layers = [
    (ACCENT2, "Layer 0 — Hardware / OS Kernel",
     "Raw disk handle (CreateFileW)  ·  WMI process events  ·  Low-level keyboard hook (WH_KEYBOARD_LL)"),
    (RGBColor(0xFF,0x80,0x00), "Layer 1 — V5 Sensor Modules",
     "NTFS Artifact Parser  ·  Entropy Map Scanner  ·  Real-time Behavior Monitor  ·  File Integrity Monitor"),
    (YELLOW, "Layer 2 — V4 Detection Engine",
     "Process Monitor (WMI + psutil)  ·  Behavioral Simulation  ·  Alert Deduplication Engine"),
    (ACCENT, "Layer 3 — AI Analysis Core",
     "Gemini 2.5 Flash Analyzer  ·  Circuit Breaker  ·  Rate Limiter  ·  Keyword Fallback Engine"),
    (GREEN,  "Layer 4 — Incident Response",
     "Emergency Snapshot Engine  ·  Incident Report Generator  ·  Evidence Chain Signer"),
    (RGBColor(0xAA,0x44,0xFF), "Layer 5 — Resilience",
     "Thread Watchdog (auto-restart)  ·  Forensic Narrative Engine  ·  Hash-Chain Logger"),
    (GREY,   "Layer 6 — Visibility",
     "Live Dashboard (FastAPI)  ·  Health API  ·  SIEM Integration (Splunk / Elastic / Syslog)"),
]

yy = 1.45
for color, title, desc in layers:
    bar = sl.shapes.add_shape(1, Inches(0.4), Inches(yy), Inches(0.12), Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    txt(sl, title, 0.65, yy+0.03, 3.2, 0.3, size=12, bold=True, color=color)
    txt(sl, desc,  3.95, yy+0.07, 9.0, 0.4, size=11, color=GREY)
    yy += 0.72

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — NTFS Parser
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(SLIDE_LAYOUT)
bg(sl)
accent_bar(sl, 0.0, 0.07, ACCENT2)
section_header(sl, "V5 MODULE 1", ACCENT2)

txt(sl, "NTFS Artifact Parser — Catching Forgeries on the Raw Disk",
    0.5, 0.62, 12.5, 0.55, size=22, bold=True, color=WHITE)
divider(sl, 1.28, ACCENT2)

# Left — What it reads
box(sl, 0.4, 1.38, 5.6, 5.8, CARD)
txt(sl, "What It Reads", 0.6, 1.5, 5.2, 0.38, size=15, bold=True, color=ACCENT2)

ntfs_items = [
    ("$MFT — Master File Table",
     "Every file record on the disk. Contains\ntwo independent timestamp sets."),
    ("$STANDARD_INFO ($SI) Timestamps",
     "The timestamps visible in Windows Explorer.\nCan be changed by any tool via the Windows API."),
    ("$FILE_NAME ($FN) Timestamps",
     "Written directly by the kernel during\nfile creation. CANNOT be altered via API."),
    ("$USN Change Journal",
     "Indestructible log of every file operation.\nShadowNet reads this to detect ghost writes."),
    ("$LogFile (Transaction Log)",
     "Captures in-progress NTFS transactions.\nCorrupted patterns signal aggressive tampering."),
]
yy = 1.95
for title, desc in ntfs_items:
    txt(sl, "◆ " + title, 0.6, yy, 5.1, 0.3, size=12, bold=True, color=ACCENT2)
    txt(sl, desc, 0.78, yy+0.28, 4.9, 0.45, size=11, color=GREY)
    yy += 0.82

# Right — Detection Logic
box(sl, 6.3, 1.38, 6.65, 5.8, CARD)
txt(sl, "Timestomp Detection Logic", 6.5, 1.5, 6.3, 0.38, size=15, bold=True, color=ACCENT)

steps = [
    ("Step 1 — Open Raw Volume",
     "ShadowNet calls CreateFileW(\"\\\\.\\\\C:\") to get a\nkernel-level handle, bypassing all OS file APIs."),
    ("Step 2 — Parse MFT Records",
     "Scans up to 10,000 MFT records per cycle,\nextracting both $SI and $FN timestamps per file."),
    ("Step 3 — Compare Timestamps",
     "If $SI Creation Time is more than 2 seconds\nearlier than $FN Creation Time → ANOMALY."),
    ("Step 4 — Cross-check USN Journal",
     "Checks if the file modification appears in the\nUSN log. No entry = ghost write detected."),
    ("Step 5 — Fire Alert",
     "TIMESTOMP_CONFIRMED alert fired with full\nraw timestamp delta as evidence."),
]
yy = 1.95
for title, desc in steps:
    badge(sl, str(steps.index((title,desc))+1), 6.35, yy+0.03, 0.32, 0.28, ACCENT, BG)
    txt(sl, title, 6.78, yy+0.02, 5.7, 0.28, size=12, bold=True, color=ACCENT)
    txt(sl, desc,  6.78, yy+0.3,  5.7, 0.44, size=11, color=GREY)
    yy += 0.92

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Entropy Scanner
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(SLIDE_LAYOUT)
bg(sl)
accent_bar(sl, 0.0, 0.07, GREEN)
section_header(sl, "V5 MODULE 2", GREEN)

txt(sl, "Entropy Map Scanner — Finding Hidden and Wiped Data",
    0.5, 0.62, 12.5, 0.55, size=22, bold=True, color=WHITE)
divider(sl, 1.28, GREEN)

# Entropy scale visual
box(sl, 0.4, 1.42, 12.5, 1.3, CARD)
txt(sl, "Shannon Entropy Scale (0 → 8 bits/byte)", 0.6, 1.5, 12.0, 0.3,
    size=12, bold=True, color=GREEN)

scale_items = [
    (0.5,  "0.0\nAll-Zero Wipe\n(Deleted/Wiped Data)",    ACCENT2),
    (3.0,  "3–5\nNormal Files\n(Text, Executables)",       GREY),
    (6.5,  "6–7\nImages / Media\n(Compressed Data)",       YELLOW),
    (9.5,  "~8.0\nEncrypted Data\n(Ransomware / Hidden Vol)", GREEN),
    (11.4, "Max\n=8.0 bits",                               GREY),
]
for x, label, color in scale_items:
    txt(sl, label, x, 1.82, 1.8, 0.85, size=9.5, color=color, align=PP_ALIGN.CENTER)

# Two detection boxes
for xi, (color, trigger, title, desc) in enumerate([
    (GREEN, "Entropy > 7.85 for 50,000+ clusters",
     "Hidden Volume / Ransomware Detected",
     "When a massive contiguous run of disk clusters all measure near-perfect mathematical randomness\n"
     "(entropy > 7.85), ShadowNet concludes that either:\n"
     "  ▸  An encrypted hidden volume exists here (like VeraCrypt), OR\n"
     "  ▸  Ransomware is actively encrypting your data in real time.\n"
     "Alert fired: HIDDEN_VOLUME_CONFIRMED"),
    (ACCENT2, "Entropy < 0.1 for large cluster runs",
     "Data Wiper / Secure Delete Detected",
     "When a large region of disk measures near-zero entropy (all bytes are identical, usually 0x00),\n"
     "ShadowNet concludes a data wiper or secure-delete tool has overwritten that region.\n"
     "  ▸  SDelete, BleachBit, or a custom wiper leaves this exact signature.\n"
     "  ▸  A Kolmogorov-Smirnov statistical test confirms the uniformity.\n"
     "Alert fired: WIPE_PATTERN_DETECTED"),
]):
    bx = xi * 6.35 + 0.4
    box(sl, bx, 2.88, 6.1, 4.25, CARD)
    s = sl.shapes.add_shape(1, Inches(bx), Inches(2.88), Inches(6.1), Inches(0.13))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    txt(sl, f"TRIGGER: {trigger}", bx+0.15, 3.07, 5.8, 0.3, size=11, bold=True, color=color)
    txt(sl, title, bx+0.15, 3.42, 5.8, 0.35, size=14, bold=True, color=WHITE)
    txt(sl, desc,  bx+0.15, 3.85, 5.8, 3.1,  size=11, color=GREY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Behavioral Monitor
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(SLIDE_LAYOUT)
bg(sl)
accent_bar(sl, 0.0, 0.07, YELLOW)
section_header(sl, "V5 MODULE 3", YELLOW)

txt(sl, "Real-time Behavioral Monitor — Human vs. Bot Detection",
    0.5, 0.62, 12.5, 0.55, size=22, bold=True, color=WHITE)
divider(sl, 1.28, YELLOW)

# Left: How it works
box(sl, 0.4, 1.42, 6.0, 5.75, CARD)
txt(sl, "Detection Pipeline", 0.6, 1.55, 5.6, 0.35, size=15, bold=True, color=YELLOW)

pipeline = [
    ("WH_KEYBOARD_LL Hook",
     "A Windows low-level keyboard hook intercepts\nevery key press system-wide at kernel message-pump level.\nEach timestamp is recorded in nanoseconds."),
    ("IKI Calculation",
     "Inter-Keystroke Interval (time between key presses)\nis extracted in milliseconds using np.diff()."),
    ("Statistical Analysis",
     "4 metrics computed:\n  ▸ Mean IKI (avg speed)\n  ▸ StDev (variation)\n  ▸ CV = StDev/Mean (relative consistency)\n  ▸ Kurtosis (distribution sharpness)"),
    ("GMM Bimodal Test",
     "A Gaussian Mixture Model checks if the typing\npattern is bimodal (human — two speed modes)\nor unimodal (bot — single mechanical pace)."),
    ("Verdict Rules",
     "BOT_CONFIRMED if 3+ rules trigger:\n  ▸ CV < 0.15 (too consistent)\n  ▸ StdDev < 20ms (superhuman)\n  ▸ Mean < 30ms (physically impossible)\n  ▸ High kurtosis burst"),
]
yy = 1.98
for step, (title, desc) in enumerate(pipeline, 1):
    badge(sl, str(step), 0.45, yy+0.03, 0.3, 0.27, YELLOW, BG)
    txt(sl, title, 0.85, yy+0.01, 5.1, 0.27, size=12, bold=True, color=YELLOW)
    txt(sl, desc,  0.85, yy+0.28, 5.1, 0.65, size=10.5, color=GREY)
    yy += 1.0

# Right: Comparison chart  
box(sl, 6.7, 1.42, 6.2, 5.75, CARD)
txt(sl, "Human vs. Bot — Statistical Signature", 6.9, 1.55, 5.8, 0.35,
    size=15, bold=True, color=YELLOW)

comparisons = [
    ("Typing Speed (Mean IKI)",  "80–200 ms",   "10–40 ms",    True),
    ("Variation (StDev)",        "30–80 ms",    "< 5 ms",      True),
    ("Pattern Shape",            "Bimodal",     "Unimodal",    True),
    ("Coeff. of Variation",      "> 0.25",      "< 0.10",      True),
    ("Kurtosis",                 "Low / Normal","Very High",   True),
    ("Physically Possible?",     "Yes",         "Often NO",    True),
]
yy = 2.02
txt(sl, "Metric",           6.9, yy, 2.2,  0.28, size=11, bold=True, color=GREY)
txt(sl, "Human",            9.1, yy, 1.65, 0.28, size=11, bold=True, color=GREEN)
txt(sl, "Automated Bot",   10.75, yy, 2.0, 0.28, size=11, bold=True, color=ACCENT2)
yy += 0.35
for metric, human, bot, flag in comparisons:
    box(sl, 6.75, yy, 5.9, 0.38, RGBColor(0x18,0x1A,0x28))
    txt(sl, metric, 6.9,  yy+0.06, 2.15, 0.28, size=10.5, color=WHITE)
    txt(sl, human,  9.1,  yy+0.06, 1.6,  0.28, size=10.5, color=GREEN)
    txt(sl, bot,   10.75, yy+0.06, 2.0,  0.28, size=10.5, color=ACCENT2, bold=True)
    yy += 0.5

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Process Monitor + FIM
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(SLIDE_LAYOUT)
bg(sl)
accent_bar(sl, 0.0, 0.07, ACCENT)
section_header(sl, "V5 MODULES 4 & 5", ACCENT)

txt(sl, "Process Monitor + File Integrity Monitor",
    0.5, 0.62, 12.5, 0.55, size=26, bold=True, color=WHITE)
divider(sl, 1.28, ACCENT)

# Process Monitor
box(sl, 0.4, 1.42, 6.2, 5.72, CARD)
s = sl.shapes.add_shape(1, Inches(0.4), Inches(1.42), Inches(6.2), Inches(0.16))
s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

txt(sl, "V4 Process Monitor — Dual Interception", 0.6, 1.65, 5.8, 0.35,
    size=15, bold=True, color=ACCENT)

pm_items = [
    ("WMI Subscription (Instant)",
     "Subscribes to Win32_Process creation events via Windows Management\n"
     "Instrumentation. The OS notifies ShadowNet the instant a process starts.\n"
     "Timeout: 500ms max wait per event cycle."),
    ("psutil Polling (Backup, 10ms)",
     "A second thread independently polls all running PIDs every 10ms\n"
     "as a fallback in case WMI misses a short-lived process."),
    ("Command-Line Keyword Matching",
     "The full command-line string is matched against 40+ suspicious keywords\n"
     "from config (vssadmin, mimikatz, wevtutil, certutil, -Enc, etc.)."),
    ("AI Analysis & Report",
     "A match triggers: Emergency Snapshot → Gemini AI analysis →\n"
     "Markdown incident report → SIEM forwarding."),
]
yy = 2.1
for title, desc in pm_items:
    txt(sl, f"▸  {title}", 0.6, yy, 5.7, 0.3, size=12, bold=True, color=ACCENT)
    txt(sl, desc, 0.75, yy+0.28, 5.5, 0.62, size=11, color=GREY)
    yy += 1.05

# FIM
box(sl, 6.85, 1.42, 6.1, 5.72, CARD)
s = sl.shapes.add_shape(1, Inches(6.85), Inches(1.42), Inches(6.1), Inches(0.16))
s.fill.solid(); s.fill.fore_color.rgb = GREEN; s.line.fill.background()

txt(sl, "File Integrity Monitor (FIM) — Real-Time File Watch", 7.05, 1.65, 5.7, 0.35,
    size=15, bold=True, color=GREEN)

fim_items = [
    ("OS Hook (Not Polling)",
     "Uses the 'watchdog' library to install an OS-level filesystem\n"
     "notification hook. No delay — events arrive in real time."),
    ("Monitored Directory",
     "Configurable path (e.g. C:/Users/Desktop). ShadowNet watches\n"
     "the entire directory tree recursively."),
    ("Suspicious Extensions",
     "Any .exe, .dll, .ps1, .bat, or .sh file being Created, Modified,\n"
     "or Moved triggers an immediate FIM alert."),
    ("Always Flag Deletions",
     "ANY file being deleted — regardless of type — is always flagged.\n"
     "Mass deletion is a primary signature of wiper malware."),
    ("No Admin Required",
     "FIM works at normal user privilege. Unlike NTFS and Entropy scanners,\n"
     "it does not need raw disk access."),
]
yy = 2.1
for title, desc in fim_items:
    txt(sl, f"▸  {title}", 7.05, yy, 5.6, 0.3, size=12, bold=True, color=GREEN)
    txt(sl, desc, 7.2, yy+0.28, 5.4, 0.62, size=11, color=GREY)
    yy += 1.0

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — AI Analysis Core
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(SLIDE_LAYOUT)
bg(sl)
accent_bar(sl, 0.0, 0.07, RGBColor(0xAA,0x44,0xFF))
section_header(sl, "AI CORE", RGBColor(0xAA,0x44,0xFF))

txt(sl, "Gemini AI Analysis Engine — Intent Over Signature",
    0.5, 0.62, 12.5, 0.55, size=26, bold=True, color=WHITE)
divider(sl, 1.28, RGBColor(0xAA,0x44,0xFF))

# Flow boxes
flow = [
    (ACCENT,                  "Suspicious\nCommand", 0.35),
    (RGBColor(0x44,0x44,0x88),"Rate\nLimiter\n20 RPM", 2.45),
    (RGBColor(0x44,0x44,0x88),"Circuit\nBreaker\n3-State", 4.55),
    (RGBColor(0x44,0x44,0x88),"Cache\nCheck\n24hr TTL", 6.65),
    (RGBColor(0xAA,0x44,0xFF),"Gemini\n2.5 Flash\nAI", 8.75),
    (GREEN,                   "Structured\nVerdict\nJSON", 10.85),
]
for color, label, x in flow:
    box(sl, x, 1.45, 1.85, 1.45, color)
    txt(sl, label, x+0.05, 1.52, 1.75, 1.3, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if x < 10.85:
        txt(sl, "→", x+1.88, 1.95, 0.5, 0.55, size=20, color=GREY, align=PP_ALIGN.CENTER)

# Fallback
box(sl, 3.5, 3.1, 6.35, 0.72, CARD)
txt(sl, "⚡  Fallback Engine (when AI unavailable)",
    3.7, 3.18, 6.0, 0.3, size=13, bold=True, color=ACCENT2)
txt(sl, "If Rate Limited or Circuit OPEN → Keyword Fallback Engine fires automatically. "
    "V5 anomalies (NTFS/Entropy) always get HIGH severity at 90% confidence.",
    3.7, 3.5, 5.95, 0.3, size=10.5, color=GREY)

# Gemini output fields
box(sl, 0.4, 4.0, 12.5, 3.2, CARD)
txt(sl, "What Gemini Returns — Structured JSON Verdict", 0.6, 4.1, 12.0, 0.35,
    size=14, bold=True, color=RGBColor(0xAA,0x44,0xFF))

fields = [
    ("is_anti_forensics", "true / false", "Is this command attempt a forensic attack?"),
    ("severity",          "CRITICAL / HIGH / MEDIUM / LOW", "How dangerous is this?"),
    ("category",          "e.g. anti_forensics, credential_theft", "Attack category"),
    ("confidence",        "0.0 → 1.0 float", "How confident is the AI in its verdict?"),
    ("mitre_attack_ttps", "e.g. T1070.001, T1003", "MITRE ATT&CK technique codes"),
    ("explanation",       "Plain English string", "What the command is doing and why it's suspicious"),
    ("recommended_action","immediate_containment / monitor / review", "Suggested next step"),
]
yy = 4.55
for field, ftype, meaning in fields:
    txt(sl, field,   0.6,  yy, 3.0,  0.3, size=11, bold=True, color=ACCENT)
    txt(sl, ftype,   3.7,  yy, 3.5,  0.3, size=11, color=YELLOW)
    txt(sl, meaning, 7.3,  yy, 5.5,  0.3, size=11, color=GREY)
    yy += 0.38

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Incident Response Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(SLIDE_LAYOUT)
bg(sl)
accent_bar(sl, 0.0, 0.07, ACCENT2)
section_header(sl, "INCIDENT RESPONSE", ACCENT2)

txt(sl, "End-to-End: From Detection to Report in Seconds",
    0.5, 0.62, 12.5, 0.55, size=26, bold=True, color=WHITE)
divider(sl, 1.28, ACCENT2)

# Timeline flow
steps = [
    (ACCENT2, "T+0ms\nDETECTION",
     "Process Monitor, NTFS Parser, Entropy Scanner, or\nFIM fires. Alert sent to deduplication layer."),
    (YELLOW,  "T+1ms\nSNAPSHOT",
     "Emergency Snapshot Engine fires 10 parallel threads\ncapturing: process tree, event logs, network connections."),
    (ACCENT,  "T+50ms\nAI ANALYSIS",
     "Command or anomaly evidence sent to Gemini 2.5 Flash.\nStructured JSON verdict returned in ~2 seconds."),
    (GREEN,   "T+3s\nREPORT",
     "Full Markdown incident report written to:\nevidence/incidents/<ID>/incident.md"),
    (RGBColor(0xAA,0x44,0xFF), "T+5s\nCHAIN SIGN",
     "SHA-256 hash of every file in the snapshot\nwritten to manifest.json for tamper-proof evidence."),
    (GREY,    "T+60s\nNARRATIVE",
     "Forensic Narrative Engine synthesizes ALL findings\ninto a unified attack story using Gemini AI."),
]
xx = 0.35
for color, title, desc in steps:
    box(sl, xx, 1.45, 2.0, 2.2, CARD)
    s = sl.shapes.add_shape(1, Inches(xx), Inches(1.45), Inches(2.0), Inches(0.15))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    txt(sl, title, xx+0.08, 1.66, 1.84, 0.55, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(sl, desc,  xx+0.1,  2.25, 1.82, 1.35, size=10,  color=GREY)
    if xx < 10.35:
        txt(sl, "→", xx+2.02, 2.35, 0.28, 0.45, size=16, color=GREY, align=PP_ALIGN.CENTER)
    xx += 2.18

# Evidence output
box(sl, 0.4, 3.85, 12.5, 3.2, CARD)
txt(sl, "Evidence Output — What Gets Saved", 0.6, 3.95, 12.0, 0.35,
    size=14, bold=True, color=ACCENT)

outputs = [
    ("evidence/incidents/<ID>/incident.md",
     "Full Markdown report: severity, MITRE codes, AI explanation, evidence list, recommended actions."),
    ("evidence/emergency_snapshots/<ID>/",
     "Folder containing: process_tree.json, event_logs.txt, network_state.json, metadata.json, manifest.json"),
    ("evidence/logs/forensic_findings.jsonl",
     "SHA-256 hash-chain log of all FORENSIC events. Tampering with any entry breaks the chain."),
    ("evidence/artifacts/ntfs_inconsistencies.jsonl",
     "Raw timestomp evidence: file path, $SI timestamp, $FN timestamp, delta in seconds."),
    ("evidence/reports/forensic_narrative_<TS>.md",
     "AI-synthesized unified attack narrative with MITRE timeline, corroboration map, evidence gaps."),
]
yy = 4.38
for path, desc in outputs:
    txt(sl, path, 0.6, yy, 5.5, 0.28, size=10.5, bold=True, color=YELLOW, italic=True)
    txt(sl, desc, 6.2, yy, 6.6, 0.28, size=10.5, color=GREY)
    yy += 0.48

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Resilience & Self-Healing
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(SLIDE_LAYOUT)
bg(sl)
accent_bar(sl, 0.0, 0.07, GREEN)
section_header(sl, "RESILIENCE", GREEN)

txt(sl, "Self-Healing System — It Doesn't Just Monitor, It Survives",
    0.5, 0.62, 12.5, 0.55, size=22, bold=True, color=WHITE)
divider(sl, 1.28, GREEN)

bullet_card(sl, 0.4, 1.42, 6.0, 3.0, "Thread Watchdog (core/watchdog.py)", [
    "All 4 sensor modules run as independent threads registered with the Watchdog.",
    "Watchdog checks heartbeats every 5 seconds.",
    "Dead or timed-out thread → factory called → fresh thread started.",
    "Crash-loop detection: thread dies < 10s after restart → +2 restart penalty.",
    "Exponential backoff: 2s → 4s → 8s ... up to 120s max.",
    "After 10 restarts → module marked permanently dead, CRITICAL log fired.",
], GREEN, 12)

bullet_card(sl, 0.4, 4.55, 6.0, 2.8, "Alert Deduplication (core/alert_deduplication.py)", [
    "SHA-256 fingerprint computed from alert type + key data fields.",
    "If identical fingerprint seen in last 5 minutes → alert suppressed.",
    "Prevents hundreds of identical incidents from a single event.",
    "LRU cache capped at 1,000 fingerprints. Oldest removed when full.",
], GREEN, 12)

bullet_card(sl, 6.65, 1.42, 6.3, 2.45, "Circuit Breaker (3 States)", [
    "CLOSED → Normal. All API requests pass through.",
    "OPEN → 5+ consecutive failures. API bypassed. Keyword Fallback active.",
    "HALF-OPEN → After 5min timeout. One test request sent to API.",
    "Success → returns to CLOSED. Failure → reopens for another 5 minutes.",
], ACCENT, 12)

bullet_card(sl, 6.65, 4.0, 6.3, 3.35, "Hash-Chain Tamper Detection", [
    "Every FORENSIC log entry hashes the previous entry (SHA-256).",
    "Forms a blockchain-style chain through all log entries.",
    "If any log entry is modified or deleted retroactively,",
    "the chain breaks at that exact sequence number.",
    "Verified at every startup via verify_log_chain().",
    "Also verified: every evidence snapshot via manifest.json comparison.",
], YELLOW, 12)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Key Claims + Live Dashboard
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(SLIDE_LAYOUT)
bg(sl)
accent_bar(sl, 0.0, 0.07, ACCENT)
accent_bar(sl, 7.41, 0.09, ACCENT2)
section_header(sl, "CAPABILITIES", ACCENT)

txt(sl, "What Makes ShadowNet v5.0 Unique",
    0.5, 0.62, 12.5, 0.55, size=26, bold=True, color=WHITE)
divider(sl, 1.28, ACCENT)

claims = [
    (ACCENT,  "Live $MFT Parsing",
     "Parses raw NTFS binary structures without a disk image, on a live running system."),
    (GREEN,   "Entropy K-S Testing",
     "Cluster-level statistical uniformity test (Kolmogorov-Smirnov) for wipe/ransomware detection."),
    (YELLOW,  "GMM Bot Classifier",
     "Gaussian Mixture Model bimodal test distinguishes human from automated input."),
    (ACCENT2, "LLM Narrative Synthesis",
     "Correlates NTFS + Entropy + Behavioral + Process findings into a single AI-authored attack story."),
    (RGBColor(0xAA,0x44,0xFF), "Triple-Fallback AI",
     "Circuit Breaker → Rate Limiter → Keyword Engine. Reports are ALWAYS generated."),
    (RGBColor(0xFF,0x80,0x00), "Crash-Loop Watchdog",
     "Detects crash-looping threads within 10s and applies exponential penalty backoff."),
]

xx = 0.38
yy = 1.48
for i, (color, title, desc) in enumerate(claims):
    box(sl, xx, yy, 4.05, 1.6, CARD)
    s = sl.shapes.add_shape(1, Inches(xx), Inches(yy), Inches(4.05), Inches(0.12))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    txt(sl, title, xx+0.15, yy+0.18, 3.75, 0.35, size=13, bold=True, color=color)
    txt(sl, desc,  xx+0.15, yy+0.56, 3.75, 0.92, size=11, color=GREY)
    if (i+1) % 3 == 0:
        xx = 0.38
        yy += 1.82
    else:
        xx += 4.3

# Bottom bar
box(sl, 0.4, 5.55, 12.5, 1.6, CARD)
txt(sl, "Technology Stack",
    0.6, 5.63, 12.0, 0.3, size=13, bold=True, color=ACCENT)
stack = [
    ("Python 3.10+", GREY), ("Google Gemini 2.5 Flash", ACCENT),
    ("ctypes / WinAPI", YELLOW), ("numpy / scipy", GREEN),
    ("scikit-learn (GMM)", GREEN), ("FastAPI + uvicorn", ACCENT),
    ("psutil + WMI", GREY), ("watchdog (FIM)", GREY),
    ("python-dotenv", GREY), ("SHA-256 / hashlib", YELLOW),
]
xx2 = 0.5
for tech, color in stack:
    badge(sl, tech, xx2, 6.05, len(tech)*0.095+0.3, 0.32, CARD, color)
    xx2 += len(tech)*0.095 + 0.45

# ─── Save ──────────────────────────────────────────────────────────────────────
out = "ShadowNet_Nexus_v5_Presentation.pptx"
prs.save(out)
print(f"Presentation saved: {out}")
print(f"Slides: 12")
